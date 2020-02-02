import pptx
import copy
from pptx.opc.constants import CONTENT_TYPE as CT, RELATIONSHIP_TYPE as RT
from pptx.parts.slide import SlidePart
from pptx.opc.constants import RELATIONSHIP_TYPE as RT
from pptx.opc.packuri import PackURI
from pptx.slide import Slide
from pptx.util import Cm as cm
import pptx.virtual as vr
import gc

prsnt = pptx.Presentation()
prsnt2 = pptx.Presentation()

slide = prsnt.slides.new_slide(prsnt.slide_layouts[0])
a = slide.shapes._spTree

slide.shapes.clear_all()
prsnt.slides.append_slide(slide)
prsnt.slides.append_slide(slide)
slide.shapes.add_table(2,2,cm(5),cm(5),cm(1),cm(1))

# slide2 = copy.deepcopy(slide)
new_slide = slide.copy()

new_slide.shapes.add_table(2,2,cm(10),cm(10),cm(1),cm(1))

prsnt.slides.append_slide(new_slide)
prsnt.slides.insert_slide(0,new_slide)
prsnt.slides.append_slide(slide)
prsnt2.slides.append_slide(slide)
prsnt2.slides.append_slide(new_slide)
prsnt.slides.remove_slide(slide)

prsnt.save('kkk.pptx')
prsnt2.save('lll.pptx')