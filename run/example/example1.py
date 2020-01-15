import pptx
import openpyexcel as excel
import math
import pathlib
from pptx.util import Cm as cm

exl_book = excel.load_workbook('C:/Users/dingo/OneDrive/2020/20191230 박선우 전시/모델 아카이빙.xlsx')
exl_sheet = exl_book.worksheets[0]
exl_rows = list(exl_sheet.rows)
exl_cols = list(exl_sheet.columns)
exl_dict = {}
for i,col in enumerate(exl_cols):
    exl_dict[col[1].value] = [c.value for c in col[2:]]
n = len(list(exl_dict.values())[0])

p = pptx.Presentation()
layout = p.slide_layouts[0]

a3_size = 42.0,29.7
p.slide_width = cm(a3_size[0])
p.slide_height = cm(a3_size[1])

tables = []
table = p.slides.new_slide(layout).shapes.new_table(6,3,cm(0),cm(0),width=cm(a3_size[0]/3), height=cm(a3_size[0]/3))

r1 = 0.5
r2 = 5
s = r1*5+r2
h1 = (r1/s)*(a3_size[1]/3)
h2 = (r2/s)*(a3_size[1]/3)
rows = list(table.table.rows)
rows[4].cells[0].merge(rows[4].cells[2])
rows[5].cells[0].merge(rows[5].cells[2])

for c in table.table.iter_cells():
    c.margin = cm(0.05)
    run = c.text_frame.paragraphs[0].add_run()
    run.font.size = cm(0.1)
    run.text = ' '
    c.fill.solid(255,255,255)
for r in rows[:5]:
    r.height = cm(h1)
rows[5].height = cm(h2)


for i, vs in enumerate(zip(*exl_dict.values())):
    # idx, room, kind, file, cr, creator, lesson, year, quality, size, case, relate, note, uid = i
    # formating values
    vs = list(vs[:-1])
    file_path = vs.pop(3)
    if file_path is None:
        file_path = ' '
    else:
        file_path = f'...{file_path[len(file_path)-20:]}'

    vs.append(file_path)
    vs[1] = {1:'209호',2:'복도',3:'3층',4:'카브 조교실'}[int(vs[1])]
    vs[2] = {'101': '구조개념',
             '201':'계단',
             '202':'구조디테일',
             '203':'트리하우스',
             '301':'전망대',
             '401':'빌딩',
             '402':'타워',
             '501':'체육관',
             '502':'지붕구조',
             '601':'교량',
             '900':'박선우 작품',
             None: ' '}[vs[2]]
    vs[3] = {'0':'자가디자인', '1':'외부디자인',None:' '}[vs[3]]
    vs[5] = {'1':'구조의 이해',
             '2':'구조시스템',
             '3':'구조디자인',
             '4':'기술스튜디오1',
             '5':'기술스튜디오2',
             '6':'구조물계획',
             None:' '}[vs[5]]
    vs[9] = {'0':'케이스 없음', '1':'케이스 있음',None:' '}[vs[9]]
    if vs[10] is not None:
        vs[10] = f'#{vs[10]}과 연관'

    # appending slide
    row_idx, col_idx = math.floor(i%9/3), i%3
    if i%9 == 0:
        slide = p.slides.add_slide(layout)
        slide.shapes.clear_all()
    else:
        pass

    # positioning table
    table.top = cm(row_idx * (a3_size[1] / 3))
    table.left = cm(col_idx * (a3_size[0] / 3))

    # texting
    for ii,(text,cell) in enumerate(zip(vs,table.table.iter_cells())):
        if (text is None or text == ' ') and ii in (2,3,4,5):
            text = ' '
            cell.fill.solid(200,0,0,0.5)
        else:
            cell.fill.solid(200,200,200,1)
        run = cell.text_frame.paragraphs[0].runs[0]
        run.font.fill.solid(0,0,0)
        run.font.size = cm(0.4)
        run.text = str(text)

    # cell as a picture frame
    pf_cell = table.table.cell(5,0)
    # putting image
    file_path = f'C:/Users/dingo/OneDrive/2020/20191230 박선우 전시/모델 아카이빙 이미지/{vs[0]}.jpg'
    if pathlib.Path(file_path).is_file():
        picture = slide.shapes.add_picture(file_path, *pf_cell.coordinate(0))
        if picture.width.cm > picture.height.cm:
            r = pf_cell.height/picture.height
            picture.height = cm(picture.height.cm*r)
            picture.width = cm(picture.width.cm*r)
        else:
            r = pf_cell.height/picture.width
            picture.height = cm(picture.height.cm*r)
            picture.width = cm(picture.width.cm*r)
            picture.rotation = -90
        picture.orient(*pf_cell.coordinate(4), 4)
    slide.shapes.append_table(table)

p.save('모형 정리.pptx')