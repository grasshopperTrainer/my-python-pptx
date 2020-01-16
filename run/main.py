import pptx
import openpyexcel as excel
from pptx.util import Cm as cm
from pptx.oxml.xmlchemy import serialize_for_reading
import math
import pathlib

p = pptx.Presentation()
layout = p.slide_layouts[0]
slide = p.slides.add_slide(layout)
table_frame = slide.shapes.add_table(rows=3, cols=3, left=0, top=0, width=cm(3), height=cm(3))
table_frame2 = slide.shapes.add_table(rows=2, cols=2, left=0, top=0, width=cm(2), height=cm(2))
table = table_frame.table
table2= table_frame2.table

table.add_row(height=cm(5))
table.add_column(width=cm(10))
table_frame.orient(posx=cm(2),posy=cm(2),ref=1)

table.delete_row()
table.delete_column()

table.join_table(table2,trim=False,pos=1)

# testing flipped access to cell object
a = table.rows[1].cells
b = table.columns[0].cells
print(a[0] == b[1])
c = table.columns[1].cells[1]

from pptx.virtual.table import Table, Slide

t = Table(rows = 3, cols= 3, left=cm(0), top=cm(0), width=cm(10), height=cm(10))
s = Slide()
p = pptx.Presentation()
s.append_in(p)
p.save('testt.pptx')